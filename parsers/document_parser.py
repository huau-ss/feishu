"""
文档解析器模块
支持 PDF, Word, Excel, PPT, TXT, HTML, XML, ZIP, CSV, RTF, 视频音频等格式
"""
import io
import logging
from abc import ABC, abstractmethod
from pathlib import Path
from typing import Optional, List, Dict, Tuple
import re

logger = logging.getLogger(__name__)


class BaseParser(ABC):
    """解析器基类"""

    @abstractmethod
    def extract(self, file_path: str | Path) -> str:
        """提取文本内容"""
        pass

    @property
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        """支持的扩展名"""
        pass

    def is_supported(self, file_path: str | Path) -> bool:
        """检查是否支持该文件"""
        ext = Path(file_path).suffix.lower()
        return ext in self.supported_extensions


class PDFParser(BaseParser):
    """PDF 解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pdf"]

    def extract(self, file_path: str | Path) -> str:
        try:
            import pymupdf
        except ImportError:
            raise ImportError("请安装 pymupdf: pip install pymupdf")

        text_parts = []
        with pymupdf.open(file_path) as doc:
            for page_num, page in enumerate(doc, 1):
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"[第{page_num}页]\n{text}")

        return "\n\n".join(text_parts)

    def extract_with_metadata(self, file_path: str | Path) -> Tuple[str, Dict]:
        """
        提取文本内容并返回元数据

        Returns:
            (文本内容, 元数据字典)
        """
        try:
            import pymupdf
        except ImportError:
            raise ImportError("请安装 pymupdf: pip install pymupdf")

        text_parts = []
        metadata: Dict = {}

        with pymupdf.open(file_path) as doc:
            metadata['page_count'] = len(doc)
            if doc.metadata:
                pdf_meta = doc.metadata
                metadata['pdf_title'] = pdf_meta.get('title', '')
                metadata['pdf_author'] = pdf_meta.get('author', '')
                metadata['pdf_subject'] = pdf_meta.get('subject', '')

            for page_num, page in enumerate(doc, 1):
                text = page.get_text("text")
                if text.strip():
                    text_parts.append(f"[第{page_num}页]\n{text}")

        return "\n\n".join(text_parts), metadata


class DocxParser(BaseParser):
    """Word 文档解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".docx", ".doc"]

    def extract(self, file_path: str | Path) -> str:
        try:
            from docx import Document
        except ImportError:
            raise ImportError("请安装 python-docx: pip install python-docx")

        doc = Document(file_path)
        paragraphs = []

        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)

        # 提取表格
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    paragraphs.append(" | ".join(cells))

        return "\n\n".join(paragraphs)


class ExcelParser(BaseParser):
    """Excel 解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".xlsx", ".xls"]

    def extract(self, file_path: str | Path) -> str:
        try:
            import openpyxl
        except ImportError:
            raise ImportError("请安装 openpyxl: pip install openpyxl")

        text_parts = []
        wb = openpyxl.load_workbook(file_path, data_only=True)

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            text_parts.append(f"[工作表: {sheet_name}]")

            for row in sheet.iter_rows(values_only=True):
                cells = [str(cell) if cell is not None else "" for cell in row]
                row_text = " | ".join(cells).strip()
                if row_text:
                    text_parts.append(row_text)

            text_parts.append("")

        return "\n".join(text_parts)


class PPTParser(BaseParser):
    """PPT 解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    def extract(self, file_path: str | Path) -> str:
        try:
            from pptx import Presentation
        except ImportError:
            raise ImportError("请安装 python-pptx: pip install python-pptx")

        prs = Presentation(file_path)
        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"[幻灯片 {slide_num}]"]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_texts.append(shape.text.strip())

            if len(slide_texts) > 1:
                text_parts.append("\n".join(slide_texts))

        return "\n\n".join(text_parts)


class HTMLParser(BaseParser):
    """HTML 解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".html", ".htm"]

    def extract(self, file_path: str | Path) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4")

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            soup = BeautifulSoup(f.read(), "html.parser")

        for script in soup(["script", "style"]):
            script.decompose()

        text = soup.get_text(separator="\n", strip=True)
        return text


class TXTParser(BaseParser):
    """纯文本解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".txt", ".md", ".json", ".yaml", ".yml", ".cfg", ".conf", ".log"]

    def extract(self, file_path: str | Path) -> str:
        encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    return f.read()
            except UnicodeDecodeError:
                continue

        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            return f.read()


class EmailParser(BaseParser):
    """邮件解析器 (.eml, .msg)"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".eml", ".msg"]

    def extract(self, file_path: str | Path) -> str:
        text_parts = []

        try:
            import email
            from email import policy

            with open(file_path, "rb") as f:
                msg = email.message_from_binary_file(f, policy=policy.default)

            text_parts.append(f"主题: {msg.get('Subject', '无')}")
            text_parts.append(f"发件人: {msg.get('From', '无')}")
            text_parts.append(f"收件人: {msg.get('To', '无')}")
            text_parts.append("")

            if msg.is_multipart():
                for part in msg.walk():
                    content_type = part.get_content_type()
                    if content_type == "text/plain":
                        payload = part.get_payload(decode=True)
                        if payload:
                            charset = part.get_content_charset() or "utf-8"
                            text_parts.append(payload.decode(charset, errors="ignore"))
            else:
                payload = msg.get_payload(decode=True)
                if payload:
                    charset = msg.get_content_charset() or "utf-8"
                    text_parts.append(payload.decode(charset, errors="ignore"))

        except Exception as e:
            logger.warning(f"解析邮件文件失败 {file_path}: {e}")
            return TXTParser().extract(file_path)

        return "\n".join(text_parts)


class XMLParser(BaseParser):
    """XML 文档解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".xml"]

    def extract(self, file_path: str | Path) -> str:
        try:
            from bs4 import BeautifulSoup
        except ImportError:
            raise ImportError("请安装 beautifulsoup4: pip install beautifulsoup4")

        encodings = ["utf-8", "gbk", "gb2312", "gb18030", "latin-1"]
        raw_content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    raw_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if raw_content is None:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw_content = f.read()

        soup = BeautifulSoup(raw_content, "lxml")
        text_parts = []

        root = soup.find()
        if root and root.name:
            text_parts.append(f"[XML根元素: {root.name}]")

        for elem in soup.find_all(True):
            tag_name = elem.name
            if elem.string and elem.string.strip():
                text_parts.append(f"[{tag_name}] {elem.string.strip()}")
            else:
                children = [c for c in elem.children if isinstance(c, str) and c.strip()]
                if children:
                    text_parts.append(f"[{tag_name}] {''.join(children).strip()}")

            for attr, value in elem.attrs.items():
                if isinstance(value, list):
                    value = ' '.join(str(v) for v in value)
                if value and str(value).strip():
                    text_parts.append(f"  @{attr}: {value}")

        return "\n".join(text_parts)


class ZIPParser(BaseParser):
    """ZIP 压缩包解析器 — 解压并递归解析内部文件"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".zip"]

    def extract(self, file_path: str | Path) -> str:
        import zipfile
        import tempfile
        import shutil

        text_parts = [f"[ZIP压缩包: {Path(file_path).name}]"]
        text_parts.append("")

        try:
            with zipfile.ZipFile(file_path, 'r') as zf:
                file_list = zf.namelist()
                text_parts.append(f"包含文件数量: {len(file_list)}")
                text_parts.append("")

                for entry in file_list:
                    if entry.endswith('/'):
                        text_parts.append(f"[目录] {entry}")
                        continue

                    info = zf.getinfo(entry)
                    text_parts.append(f"[文件] {entry} ({info.file_size} bytes)")
                    text_parts.append("")

                    if info.file_size == 0:
                        continue

                    safe_ext = Path(entry).suffix.lower()

                    if safe_ext in ['.txt', '.md', '.log', '.cfg', '.conf', '.json', '.yaml', '.yml', '.xml', '.html', '.htm']:
                        try:
                            raw = zf.read(entry)
                            for enc in ['utf-8', 'gbk', 'gb2312', 'latin-1']:
                                try:
                                    content = raw.decode(enc)
                                    break
                                except UnicodeDecodeError:
                                    continue
                            else:
                                content = raw.decode('utf-8', errors='ignore')

                            if safe_ext == '.xml':
                                from bs4 import BeautifulSoup
                                soup = BeautifulSoup(content, "lxml")
                                content = soup.get_text(separator="\n", strip=True)

                            lines = content.split('\n')
                            snippet = '\n'.join(lines[:100])
                            if len(lines) > 100:
                                snippet += f"\n... (共 {len(lines)} 行，内容已截断)"
                            text_parts.append(snippet)
                            text_parts.append("")
                        except Exception as e:
                            text_parts.append(f"  (解析失败: {e})")
                            text_parts.append("")

                    elif safe_ext in ['.py', '.js', '.java', '.cpp', '.c', '.go', '.rs', '.sh', '.bat', '.sql']:
                        try:
                            raw = zf.read(entry)
                            for enc in ['utf-8', 'latin-1']:
                                try:
                                    content = raw.decode(enc)
                                    break
                                except UnicodeDecodeError:
                                    continue
                            else:
                                content = raw.decode('utf-8', errors='ignore')

                            lines = content.split('\n')
                            snippet = '\n'.join(lines[:50])
                            if len(lines) > 50:
                                snippet += f"\n... (共 {len(lines)} 行，代码已截断)"
                            text_parts.append(f"[代码片段 - {entry}]")
                            text_parts.append(snippet)
                            text_parts.append("")
                        except Exception:
                            pass

                    elif safe_ext in ['.csv']:
                        try:
                            raw = zf.read(entry)
                            for enc in ['utf-8', 'gbk', 'latin-1']:
                                try:
                                    content = raw.decode(enc)
                                    break
                                except UnicodeDecodeError:
                                    continue
                            else:
                                content = raw.decode('utf-8', errors='ignore')

                            lines = content.split('\n')
                            snippet = '\n'.join(lines[:50])
                            if len(lines) > 50:
                                snippet += f"\n... (共 {len(lines)} 行，表格已截断)"
                            text_parts.append(f"[表格 - {entry}]")
                            text_parts.append(snippet)
                            text_parts.append("")
                        except Exception:
                            pass

                    elif safe_ext in ['.pdf', '.docx', '.doc', '.xlsx', '.xls', '.pptx', '.ppt', '.rtf']:
                        text_parts.append(f"  (二进制文件，需专用解析器处理)")
                        text_parts.append("")

        except zipfile.BadZipFile:
            logger.warning(f"ZIP文件损坏或格式错误: {file_path}")
            return f"[错误] 无法解压ZIP文件: {file_path}"

        return "\n".join(text_parts)


class CSVParser(BaseParser):
    """CSV 表格解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".csv"]

    def extract(self, file_path: str | Path) -> str:
        import csv

        encodings = ["utf-8", "gbk", "gb2312", "gb18030", "utf-8-sig", "latin-1"]
        raw_content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    raw_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if raw_content is None:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw_content = f.read()

        text_parts = []
        lines = raw_content.split('\n')

        dialect = None
        try:
            dialect = csv.Sniffer().sniff(raw_content[:8192])
        except Exception:
            pass

        reader = csv.reader(lines, dialect=dialect)
        headers = next(reader, None)

        if headers:
            text_parts.append("| " + " | ".join(str(h).strip() for h in headers) + " |")
            text_parts.append("| " + " | ".join("---" for _ in headers) + " |")

            for i, row in enumerate(reader):
                if i >= 1000:
                    text_parts.append(f"\n... (共超过 {i + 1} 行，数据已截断)")
                    break
                row_cells = [str(cell).strip() for cell in row]
                if any(row_cells):
                    text_parts.append("| " + " | ".join(row_cells) + " |")
        else:
            for line in lines[:200]:
                if line.strip():
                    cells = next(csv.reader([line]), [])
                    text_parts.append(" | ".join(c.strip() for c in cells if c.strip()))

        return "\n".join(text_parts)


class RTFParser(BaseParser):
    """RTF 富文本格式解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".rtf"]

    def extract(self, file_path: str | Path) -> str:
        try:
            import striprtf.striprtf as rtf
        except ImportError:
            logger.warning("striprtf 未安装，尝试替代方法解析RTF")
            return self._extract_fallback(file_path)

        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        raw_content = None

        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    raw_content = f.read()
                break
            except UnicodeDecodeError:
                continue

        if raw_content is None:
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                raw_content = f.read()

        text = rtf.rtf_to_text(raw_content)
        return text.strip() if text else ""

    def _extract_fallback(self, file_path: str | Path) -> str:
        encodings = ["utf-8", "gbk", "gb2312", "latin-1"]
        for encoding in encodings:
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    raw = f.read()
                break
            except UnicodeDecodeError:
                continue

        import re
        raw = re.sub(r'\\[a-z]+\d*\s?', ' ', raw)
        raw = re.sub(r'\{|\}', '', raw)
        raw = re.sub(r"\\'([0-9a-fA-F]{2})", lambda m: chr(int(m.group(1), 16)), raw)
        raw = re.sub(r'\n+', '\n', raw)
        return raw.strip()


class VideoParser(BaseParser):
    """视频文件解析器 — 提取音频并通过 Whisper 进行语音转文字"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".mp4", ".avi", ".mov", ".mkv", ".flv", ".wmv", ".webm", ".mp3", ".wav", ".m4a", ".flac", ".ogg"]

    def extract(self, file_path: str | Path) -> str:
        import os
        import tempfile

        video_extensions = {'.mp4', '.avi', '.mov', '.mkv', '.flv', '.wmv', '.webm'}
        audio_extensions = {'.mp3', '.wav', '.m4a', '.flac', '.ogg'}
        ext = Path(file_path).suffix.lower()

        text_parts = [f"[视频/音频文件: {Path(file_path).name}]"]
        text_parts.append("")

        file_size_mb = os.path.getsize(file_path) / (1024 * 1024)
        text_parts.append(f"文件大小: {file_size_mb:.2f} MB")
        text_parts.append("")

        try:
            import subprocess
            subprocess.run(["ffmpeg", "-version"], capture_output=True, check=True)
        except (subprocess.CalledProcessError, FileNotFoundError):
            text_parts.append("[警告] FFmpeg 未安装，无法提取音频进行语音转文字。")
            text_parts.append("请安装 FFmpeg: https://ffmpeg.org/download.html")
            text_parts.append("Linux/macOS: sudo apt install ffmpeg / brew install ffmpeg")
            text_parts.append("Windows: 从 https://ffmpeg.org/download.html 下载并添加到 PATH")
            return "\n".join(text_parts)

        try:
            import whisper
        except ImportError:
            text_parts.append("[警告] OpenAI Whisper 未安装，无法进行语音转文字。")
            text_parts.append("请安装: pip install openai-whisper")
            text_parts.append("或安装轻量版: pip install faster-whisper")
            return "\n".join(text_parts)

        with tempfile.NamedTemporaryFile(suffix='.wav', delete=False) as tmp_audio:
            tmp_audio_path = tmp_audio.name

        try:
            if ext in video_extensions:
                cmd = [
                    "ffmpeg", "-i", str(file_path),
                    "-vn", "-acodec", "pcm_s16le",
                    "-ar", "16000", "-ac", "1",
                    "-y", tmp_audio_path
                ]
            else:
                if ext == '.wav':
                    tmp_audio_path = str(file_path)
                else:
                    cmd = [
                        "ffmpeg", "-i", str(file_path),
                        "-vn", "-acodec", "pcm_s16le",
                        "-ar", "16000", "-ac", "1",
                        "-y", tmp_audio_path
                    ]

            if ext in video_extensions or ext not in {'.wav'}:
                result = subprocess.run(
                    cmd,
                    capture_output=True, text=True, timeout=300
                )
                if result.returncode != 0:
                    text_parts.append(f"[警告] 音频提取失败: {result.stderr[:200]}")
                    return "\n".join(text_parts)

            try:
                model_name = "base"
                model = whisper.load_model(model_name)
                text_parts.append(f"[Whisper模型: {model_name}]")
                text_parts.append("")

                result = model.transcribe(tmp_audio_path, language="zh", task="transcribe")
                segments = result.get("segments", [])

                if not segments:
                    text_parts.append("[转写结果]")
                    text_parts.append(result.get("text", "").strip())
                else:
                    text_parts.append(f"[转写片段数: {len(segments)}]")
                    text_parts.append("")
                    for seg in segments:
                        start = seg.get("start", 0)
                        end = seg.get("end", 0)
                        seg_text = seg.get("text", "").strip()
                        mins = int(start // 60)
                        secs = int(start % 60)
                        text_parts.append(f"[{mins:02d}:{secs:02d} - {int(end//60):02d}:{int(end%60):02d}] {seg_text}")

            except Exception as e:
                logger.warning(f"Whisper 转写失败: {e}")
                text_parts.append(f"[警告] 语音转文字失败: {e}")
                text_parts.append("提示: 可尝试安装 faster-whisper 以提高性能: pip install faster-whisper")

        finally:
            if tmp_audio_path != str(file_path) and os.path.exists(tmp_audio_path):
                try:
                    os.unlink(tmp_audio_path)
                except Exception:
                    pass

        return "\n".join(text_parts)


class ParserFactory:
    """解析器工厂"""

    def __init__(self):
        self._parsers: List[BaseParser] = [
            PDFParser(),
            DocxParser(),
            ExcelParser(),
            PPTParser(),
            HTMLParser(),
            TXTParser(),
            EmailParser(),
            XMLParser(),
            ZIPParser(),
            CSVParser(),
            RTFParser(),
            VideoParser(),
        ]

    def get_parser(self, file_path: str | Path) -> Optional[BaseParser]:
        """根据文件类型获取解析器"""
        for parser in self._parsers:
            if parser.is_supported(file_path):
                return parser
        return None

    def extract(self, file_path: str | Path) -> str:
        """提取文件内容"""
        parser = self.get_parser(file_path)
        if parser is None:
            raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")

        content = parser.extract(file_path)
        logger.info(f"成功解析文件: {file_path}, 提取文本 {len(content)} 字符")
        return content

    def extract_with_metadata(self, file_path: str | Path) -> Tuple[str, Dict]:
        """
        提取文件内容并返回元数据（仅 PDF 支持）

        Returns:
            (文本内容, 元数据字典)

        Raises:
            ValueError: 不支持的文件类型
        """
        parser = self.get_parser(file_path)
        if parser is None:
            raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")

        if isinstance(parser, PDFParser):
            content, meta = parser.extract_with_metadata(file_path)
        else:
            content = parser.extract(file_path)
            meta = {}

        logger.info(f"成功解析文件: {file_path}, 提取文本 {len(content)} 字符")
        return content, meta

    def supported_extensions(self) -> List[str]:
        """所有支持的扩展名"""
        extensions = []
        for parser in self._parsers:
            extensions.extend(parser.supported_extensions)
        return list(set(extensions))


# 全局解析器实例
_default_parser: Optional[ParserFactory] = None


def get_parser() -> ParserFactory:
    """获取全局解析器实例"""
    global _default_parser
    if _default_parser is None:
        _default_parser = ParserFactory()
    return _default_parser


def extract(file_path: str | Path) -> str:
    """便捷函数：提取文件内容"""
    return get_parser().extract(file_path)


def extract_with_metadata(file_path: str | Path) -> Tuple[str, Dict]:
    """便捷函数：提取文件内容并返回元数据（仅 PDF 支持）"""
    return get_parser().extract_with_metadata(file_path)
